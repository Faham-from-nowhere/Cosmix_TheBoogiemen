from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_deck():
    prj = Presentation()

    def add_slide(title_text, prj_obj, layout_idx=1):
        slide = prj_obj.slides.add_slide(prj_obj.slide_layouts[layout_idx])
        
        # Enforce dark background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(12, 12, 16)
        
        if slide.shapes.title:
            title = slide.shapes.title
            title.text = title_text
            # Format title
            title_tf = title.text_frame
            for p in title_tf.paragraphs:
                p.font.color.rgb = RGBColor(97, 218, 251) # Sleek blue/cyan
                p.font.name = "Arial"
                p.font.bold = True
        return slide

    def add_bullets(slide, content_list):
        if len(slide.shapes.placeholders) > 1:
            body_shape = slide.shapes.placeholders[1]
            tf = body_shape.text_frame
            tf.clear()
            for i, points in enumerate(content_list):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = points
                p.font.color.rgb = RGBColor(230, 230, 230)
                p.font.size = Pt(24)
                p.font.name = "Arial"

    # Slide 1: Title
    slide0 = add_slide("SAR Ship Detection Pipeline", prj, layout_idx=0)
    subtitle = slide0.placeholders[1]
    subtitle.text = "Cosmix The Boogiemen\nAn Explainable, Defended, Edge-Deployed Radar Target Protocol"
    for p in subtitle.text_frame.paragraphs:
        p.font.color.rgb = RGBColor(200, 200, 200)

    # Slide 2: The Challenge
    s1 = add_slide("The Challenge - Noise & Trust", prj)
    add_bullets(s1, [
        "SAR Imagery is notoriously noisy, saturated with speckle artifacts.",
        "Small target obfuscation: Target ships blend into background ocean clutter.",
        "\"Black Box\" Models: Operators cannot trust AI targeting blindly.",
        "Network Vulnerability: Standard nets collapse when adversarial noise is injected."
    ])

    # Slide 3: Phase 1 & 2
    s2 = add_slide("Phase 1 & 2: MVP & GAN Data Engine", prj)
    add_bullets(s2, [
        "YOLOv8-nano detection core dynamically processes massive raw structures.",
        "Programmatic translation of COCO JSON to normalized YOLO labels.",
        "PyTorch Deep Convolutional GAN (DCGAN) synthesizes raw 64x64 footprints.",
        "Seamless injection of synthesized chips perfectly stretches dataset fidelity."
    ])

    # Slide 4: Phase 3
    s3 = add_slide("Phase 3: Concept-Based Explainability", prj, layout_idx=5)
    tf = s3.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(1)).text_frame
    p = tf.add_paragraph()
    p.text = "Stripped away 'Black Box' nature to extract structural radar metrics (RCS & Laplacian Variance)."
    p.font.color.rgb = RGBColor(230, 230, 230)
    p.font.size = Pt(22)
    try:
        s3.shapes.add_picture("web/assets/explainable_output.jpg", Inches(2.8), Inches(2.2), height=Inches(4.5))
    except Exception as e:
        print("Missing image:", e)

    # Slide 5: Phase 4
    s4 = add_slide("Phase 4: Zero-Server WebAssembly", prj)
    add_bullets(s4, [
        "Bypassed entire GPU clouds by flattening the neural architecture.",
        "ONNX Graph: Converted best.pt into an agnostic standalone .onnx graph.",
        "WebAssembly Runtime: onnxruntime-web execution exclusively via client CPU.",
        "Sleek Glassmorphism HTML/JS frontend requiring purely zero backend logic."
    ])

    # Slide 6: Phase 5
    s5 = add_slide("Phase 5: Adversarial Defense", prj, layout_idx=5)
    tf2 = s5.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(1)).text_frame
    p2 = tf2.add_paragraph()
    p2.text = "Spatial protective shield restoring network failure via adaptive defense filtering."
    p2.font.color.rgb = RGBColor(230, 230, 230)
    p2.font.size = Pt(22)
    try:
        s5.shapes.add_picture("web/assets/attack_pred.jpg", Inches(0.5), Inches(2.2), width=Inches(4.2))
        s5.shapes.add_picture("web/assets/defend_pred.jpg", Inches(5.3), Inches(2.2), width=Inches(4.2))
    except Exception as e:
        print("Missing image:", e)

    # Slide 7: Phase 6
    s6 = add_slide("Phase 6: Sentinel Scale Inference", prj, layout_idx=5)
    tf3 = s6.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(1)).text_frame
    p3 = tf3.add_paragraph()
    p3.text = "800x800 tiling sliding windows scaling dynamically across multi-gigapixel satellite matrices."
    p3.font.color.rgb = RGBColor(230, 230, 230)
    p3.font.size = Pt(22)
    try:
        s6.shapes.add_picture("web/assets/sentinel2_final_output.jpg", Inches(1.5), Inches(2.2), width=Inches(7))
    except Exception as e:
        print("Missing image:", e)

    # Slide 8: Demo (Conclusion)
    s7 = add_slide("Mission Accomplished 🚀", prj)
    add_bullets(s7, [
        "Pipeline fully mapped, defended, and edge-deployed.",
        "Live interactive interface available on our GitHub Pages Repository.",
        "",
        "Thank You!",
        "Cosmix The Boogiemen"
    ])

    out_path = "Cosmix_TheBoogiemen_Presentation.pptx"
    prj.save(out_path)
    print(f"Successfully generated {out_path}!")

if __name__ == '__main__':
    create_deck()
